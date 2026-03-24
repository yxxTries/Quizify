import io
import pdfplumber
from pptx import Presentation


def extract_text(file_bytes: bytes, filename: str) -> str:
    """
    Extract plain text from a PDF or PPTX file.
    Returns the raw text string.
    Raises ValueError for unsupported types or empty results.
    """
    name = filename.lower()

    if name.endswith(".pdf"):
        return _extract_pdf(file_bytes)
    elif name.endswith(".pptx"):
        return _extract_pptx(file_bytes)
    else:
        raise ValueError(f"Unsupported file type: {filename}. Please upload a .pdf or .pptx file.")


def _extract_pdf(file_bytes: bytes) -> str:
    text_parts = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text.strip())

    result = "\n\n".join(text_parts)
    if not result.strip():
        raise ValueError(
            "No readable text found in the PDF. "
            "The file may be a scanned image — please use a text-based PDF."
        )
    return result


def _extract_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    text_parts = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = " ".join(run.text for run in para.runs).strip()
                if line:
                    slide_texts.append(line)

        if slide_texts:
            text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))

    result = "\n\n".join(text_parts)
    if not result.strip():
        raise ValueError(
            "No readable text found in the PowerPoint. "
            "The file may contain only images or empty slides."
        )
    return result
